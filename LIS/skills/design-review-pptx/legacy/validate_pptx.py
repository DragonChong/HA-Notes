#!/usr/bin/env python3
"""Quick validation: no float EMU in XML, all slides readable."""
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation


def main() -> int:
    path = Path(sys.argv[1])
    with zipfile.ZipFile(path) as z:
        bad = []
        for name in z.namelist():
            if not name.endswith(".xml"):
                continue
            text = z.read(name).decode("utf-8", errors="ignore")
            for m in re.finditer(r'(?:x="|y="|cy="|cx="|w="|h=")(\d+\.\d+)"', text):
                bad.append((name, m.group(1)))
        print(f"Float EMU in XML: {len(bad)}")
        if bad:
            for item in bad[:10]:
                print(" ", item)
            return 1

    prs = Presentation(str(path))
    print(f"Slides loaded: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        title = slide.shapes.title.text if slide.shapes.title else "?"
        placeholders = sum(
            1 for sh in slide.shapes if getattr(sh, "is_placeholder", False)
        )
        zero_ph = [
            sh.name
            for sh in slide.shapes
            if getattr(sh, "is_placeholder", False) and sh.width == 0
        ]
        print(f"  {i}: {title[:50]} | placeholders={placeholders} zero={zero_ph}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
