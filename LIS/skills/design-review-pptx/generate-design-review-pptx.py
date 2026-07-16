#!/usr/bin/env python3
"""
Generate HA/LIS CP3 design review PPTX from Markdown using the canonical template.

Usage:
    python generate-design-review-pptx.py <input.md> [output.pptx]
    python generate-design-review-pptx.py deck.md

If output path is omitted, writes {input-basename}.pptx next to the markdown file.
Template: ha-lis-design-review-template.pptx (same master as LIS-10672 / LIS-10583 decks).
"""

from __future__ import annotations

import re
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Content area geometry (10" x 7.5" slide; title band ~1.3")
CONTENT_LEFT = Inches(0.65)
CONTENT_WIDTH = Inches(8.7)
CONTENT_TOP = Inches(1.45)
CONTENT_BOTTOM = Inches(6.85)
ROW_HEIGHT = Inches(0.42)
TABLE_HEADER_PT = Pt(12)
TABLE_BODY_PT = Pt(11)
BODY_PT = Pt(18)
CAPTION_PT = Pt(14)
CODE_PT = Pt(10)

SKILL_DIR = Path(__file__).resolve().parent
TEMPLATE = SKILL_DIR / "ha-lis-design-review-template.pptx"

# Directory the input Markdown file lives in — relative image references
# (`![](diagram.png)`) resolve against this. Set by main() before generation.
MD_DIR = Path(".")

LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_TITLE_ONLY = 7

IMAGE_RE = re.compile(r"^!\[[^\]]*\]\(([^)]+)\)$")


def _emu(value) -> int:
    """Integer EMU — float values in xfrm attributes corrupt Office XML."""
    return int(round(value))


def _add_title_only_slide(prs: Presentation, title: str):
    """Content slide with title only (no body placeholder)."""
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
    title_shape = s.shapes.title
    title_shape.text = title
    # Title-only layout in the HA template uses oversized centered title text.
    # Reposition it to match the standard content-slide title band.
    title_shape.left = Inches(0.85)
    title_shape.top = Inches(0.14)
    title_shape.width = Inches(9.1)
    title_shape.height = Inches(0.72)
    tf = title_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.name = "Calibri Light"
            run.font.size = Pt(28)
    return s


def delete_all_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        slide_id = prs.slides._sldIdLst[0]
        r_id = slide_id.rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def body_placeholder(slide):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx != 0:
            return shape
    raise RuntimeError("Body placeholder not found")


def clear_text_frame(text_frame) -> None:
    text_frame.clear()
    if not text_frame.paragraphs:
        text_frame.add_paragraph()


def parse_slides(md_content: str) -> list[dict]:
    """Parse design-review Markdown into slide dicts."""
    parts = re.split(r"<!--\s*Slide number:\s*\d+\s*-->\s*", md_content)
    slides = []

    for part in parts:
        part = part.strip()
        if not part:
            continue

        lines = part.splitlines()
        if not lines or not lines[0].startswith("# "):
            continue

        title = lines[0][2:].strip()
        body = [line.rstrip() for line in lines[1:] if line.strip()]
        slides.append({"title": title, "body": body})

    return slides


def extract_image(body: list[str]) -> tuple[str | None, list[str]]:
    """Pull the first `![alt](path)` reference out of a slide body.

    Returns (image_ref, remaining_body). Only one image per slide is
    supported, matching the documented slide-type patterns; any further
    image lines are left in the body as text.
    """
    image_ref: str | None = None
    remaining: list[str] = []

    for line in body:
        match = IMAGE_RE.match(line.strip())
        if match and image_ref is None:
            image_ref = match.group(1)
            continue
        remaining.append(line)

    return image_ref, remaining


def is_table_line(line: str) -> bool:
    return line.strip().startswith("|")


def parse_markdown_table(body: list[str]) -> tuple[list[str], list[list[str]], list[str]]:
    """Return headers, rows, and footnotes from body lines."""
    table_lines: list[str] = []
    footnotes: list[str] = []
    in_table = False

    for line in body:
        if is_table_line(line):
            table_lines.append(line)
            in_table = True
        elif in_table:
            footnotes.append(line.strip())

    if not table_lines:
        return [], [], footnotes

    rows_raw = [
        [cell.strip() for cell in line.strip().strip("|").split("|")]
        for line in table_lines
        if not re.match(r"^\|\s*-+", line.strip())
    ]
    if not rows_raw:
        return [], [], footnotes

    return rows_raw[0], rows_raw[1:], footnotes


def parse_code_block(body: list[str]) -> tuple[str | None, list[str]]:
    """Extract fenced code block and trailing notes."""
    text = "\n".join(body)
    match = re.search(r"```(?:\w+)?\n(.*?)```", text, re.DOTALL)
    if not match:
        return None, body

    code = match.group(1).strip()
    before, after = text[: match.start()], text[match.end() :]
    notes = [line.strip() for line in (before + after).splitlines() if line.strip()]
    return code, notes


def classify_slide(slide: dict, index: int, *, has_image: bool = False) -> str:
    title = slide["title"]
    body = slide["body"]

    if title.strip().upper() == "Q&A":
        return "qa"
    if index == 0:
        # Slide order guarantees the first slide is the title slide — don't
        # gate this on the literal string "CP3" appearing in the body, since
        # that breaks silently for any other review forum or a typo.
        return "title"
    if title.strip().lower() == "agenda":
        return "agenda"
    if "```" in "\n".join(body):
        # Check for a code fence before table lines: a code sample (e.g. SQL
        # using `||` concatenation) can contain a line starting with "|"
        # and would otherwise be misclassified as a table slide.
        return "code"
    if any(is_table_line(line) for line in body):
        return "table"
    if has_image:
        caption_lines = [line.strip() for line in body if line.strip()]
        if len(caption_lines) <= 2:
            return "diagram"
    return "bullets"


def _set_cell_text(cell, text: str, *, bold: bool = False, size: Pt = TABLE_BODY_PT) -> None:
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = "Calibri"
        paragraph.font.size = size
        paragraph.font.bold = bold
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)


def add_title_slide(prs: Presentation, slide: dict) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    body = slide["body"]

    service_line = ""
    meta_lines: list[str] = []
    for line in body:
        if line.startswith("(") and line.endswith(")"):
            service_line = line
        else:
            meta_lines.append(line)

    title_text = slide["title"]
    if service_line:
        title_text = f"{title_text}\n{service_line}"

    s.shapes.title.text = title_text

    for shape in s.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = "\n".join(meta_lines)
            break


def add_agenda_slide(prs: Presentation, slide: dict) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
    s.shapes.title.text = slide["title"]
    body = body_placeholder(s)
    clear_text_frame(body.text_frame)

    for i, item in enumerate(slide["body"]):
        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.text = item
        p.font.bold = True


def add_bullet_slide(prs: Presentation, slide: dict) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_CONTENT])
    s.shapes.title.text = slide["title"]
    body = body_placeholder(s)
    clear_text_frame(body.text_frame)

    for i, line in enumerate(slide["body"]):
        level = 0
        text = line
        if line.startswith("- "):
            text = line[2:]
            level = 1
        elif line.startswith("  ") or line.startswith("\t"):
            text = line.strip()
            level = 1

        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.font.name = "Calibri"
        p.font.size = BODY_PT
        p.space_after = Pt(6)


def add_table_slide(prs: Presentation, slide: dict) -> None:
    headers, rows, footnotes = parse_markdown_table(slide["body"])
    if not headers:
        add_bullet_slide(prs, slide)
        return

    s = _add_title_only_slide(prs, slide["title"])

    table_rows = [headers, *rows]
    row_count = len(table_rows)
    col_count = len(headers)
    table_height = _emu(ROW_HEIGHT * row_count)
    footnote_space = Inches(0.55) if footnotes else 0
    available = _emu(CONTENT_BOTTOM - CONTENT_TOP - footnote_space)
    if table_height < _emu(available * 0.75):
        table_top = _emu(CONTENT_TOP + (available - table_height) / 2)
    else:
        table_top = _emu(CONTENT_TOP)
        table_height = min(table_height, available)

    table_shape = s.shapes.add_table(
        row_count,
        col_count,
        CONTENT_LEFT,
        table_top,
        CONTENT_WIDTH,
        table_height,
    )
    table = table_shape.table
    total_width = _emu(CONTENT_WIDTH)
    col_width = total_width // col_count
    remainder = total_width - col_width * col_count
    for c in range(col_count):
        table.columns[c].width = col_width + (remainder if c == col_count - 1 else 0)
    for r, row in enumerate(table_rows):
        for c, value in enumerate(row):
            _set_cell_text(
                table.cell(r, c),
                value,
                bold=(r == 0),
                size=TABLE_HEADER_PT if r == 0 else TABLE_BODY_PT,
            )

    if footnotes:
        note_top = _emu(table_top + table_height + Inches(0.12))
        box = s.shapes.add_textbox(CONTENT_LEFT, note_top, CONTENT_WIDTH, Inches(0.45))
        tf = box.text_frame
        tf.word_wrap = True
        for i, note in enumerate(footnotes):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = note
            p.font.name = "Calibri"
            p.font.size = TABLE_BODY_PT
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def add_code_slide(prs: Presentation, slide: dict) -> None:
    code, notes = parse_code_block(slide["body"])
    if not code:
        add_bullet_slide(prs, slide)
        return

    s = _add_title_only_slide(prs, slide["title"])
    code_height = Inches(min(0.22 * max(len(code.splitlines()), 1) + 0.3, 4.8))
    box = s.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, code_height)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = "Consolas"
    p.font.size = CODE_PT

    if notes:
        note_box = s.shapes.add_textbox(
            CONTENT_LEFT,
            CONTENT_TOP + code_height + Inches(0.15),
            CONTENT_WIDTH,
            Inches(0.8),
        )
        note_tf = note_box.text_frame
        note_tf.word_wrap = True
        for i, note in enumerate(notes):
            np = note_tf.paragraphs[0] if i == 0 else note_tf.add_paragraph()
            np.text = note
            np.font.name = "Calibri"
            np.font.size = TABLE_BODY_PT


def _resolve_image_path(image_ref: str) -> Path:
    image_path = Path(image_ref)
    if not image_path.is_absolute():
        image_path = (MD_DIR / image_path).resolve()
    return image_path


def add_diagram_slide(prs: Presentation, slide: dict, image_ref: str) -> None:
    """Large centred diagram with optional caption below (avoids title overlap)."""
    image_path = _resolve_image_path(image_ref)
    caption_lines = [line.strip() for line in slide["body"] if line.strip()]
    if not image_path.exists():
        print(f"Warning: image not found, skipping embed: {image_path}", file=sys.stderr)
        add_bullet_slide(prs, slide)
        return

    s = _add_title_only_slide(prs, slide["title"])

    caption_height = Inches(0.35 * len(caption_lines) + 0.15) if caption_lines else 0
    image_area_bottom = _emu(CONTENT_BOTTOM - caption_height - Inches(0.1))
    image_area_height = _emu(image_area_bottom - CONTENT_TOP)

    pic = s.shapes.add_picture(
        str(image_path),
        CONTENT_LEFT,
        CONTENT_TOP,
        width=CONTENT_WIDTH,
    )
    if pic.height > image_area_height:
        scale = image_area_height / pic.height
        pic.width = _emu(pic.width * scale)
        pic.height = _emu(pic.height * scale)
    pic.left = _emu(CONTENT_LEFT + (CONTENT_WIDTH - pic.width) / 2)
    pic.top = _emu(CONTENT_TOP + (image_area_height - pic.height) / 2)

    if caption_lines:
        cap_top = _emu(pic.top + pic.height + Inches(0.12))
        cap = s.shapes.add_textbox(CONTENT_LEFT, cap_top, CONTENT_WIDTH, caption_height)
        tf = cap.text_frame
        tf.word_wrap = True
        for i, line in enumerate(caption_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = "Calibri"
            p.font.size = CAPTION_PT
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
            p.alignment = PP_ALIGN.CENTER


def embed_image_beside_text(slide, image_ref: str, kind: str) -> None:
    """Image on the right when substantial bullet text shares the slide."""
    image_path = _resolve_image_path(image_ref)
    if not image_path.exists():
        print(f"Warning: image not found, skipping embed: {image_path}", file=sys.stderr)
        return

    if kind in ("bullets", "agenda"):
        try:
            body = body_placeholder(slide)
            body.left = CONTENT_LEFT
            body.top = CONTENT_TOP
            body.width = Inches(4.35)
            body.height = _emu(CONTENT_BOTTOM - CONTENT_TOP)
        except RuntimeError:
            pass

    slide.shapes.add_picture(
        str(image_path),
        Inches(5.25),
        CONTENT_TOP,
        width=Inches(4.1),
    )


def add_qa_slide(prs: Presentation) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
    title = s.shapes.title
    title.text = "Q&A"
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def build_presentation(slides: list[dict], output_path: Path) -> Presentation:
    if not TEMPLATE.exists():
        raise FileNotFoundError(
            f"Template not found: {TEMPLATE}\n"
            "Ensure ha-lis-design-review-template.pptx is in the skill folder."
        )

    shutil.copy(TEMPLATE, output_path)
    prs = Presentation(output_path)
    delete_all_slides(prs)

    handlers = {
        "title": add_title_slide,
        "agenda": add_agenda_slide,
        "bullets": add_bullet_slide,
        "table": add_table_slide,
        "code": add_code_slide,
        "qa": lambda p, _: add_qa_slide(p),
    }

    for i, slide in enumerate(slides):
        image_ref, body_wo_image = extract_image(slide["body"])
        working_slide = {"title": slide["title"], "body": body_wo_image}

        kind = classify_slide(working_slide, i, has_image=bool(image_ref))
        if kind == "diagram":
            add_diagram_slide(prs, working_slide, image_ref)
            continue

        handler = handlers[kind]
        if kind == "qa":
            handler(prs, working_slide)
        else:
            handler(prs, working_slide)

        if image_ref:
            embed_image_beside_text(prs.slides[-1], image_ref, kind)

    return prs


def main() -> None:
    global MD_DIR

    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    md_path = Path(sys.argv[1]).resolve()
    if not md_path.exists():
        print(f"Error: Markdown file not found: {md_path}", file=sys.stderr)
        sys.exit(1)

    MD_DIR = md_path.parent

    output_path = (
        Path(sys.argv[2]).resolve()
        if len(sys.argv) > 2
        else md_path.with_suffix(".pptx")
    )

    md_content = md_path.read_text(encoding="utf-8")
    slides = parse_slides(md_content)
    if not slides:
        print("Error: No slides parsed from markdown.", file=sys.stderr)
        sys.exit(1)

    prs = build_presentation(slides, output_path)
    prs.save(output_path)
    print(f"Created: {output_path} ({len(slides)} slides)")


if __name__ == "__main__":
    main()
